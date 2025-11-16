"""
Smart RAG Implementation with Pluggable Vision Models

Key Changes:
- Vision models are now pluggable via VisionModelFactory
- Users can select: Moondream2, Qwen3-VL, or InternVL3.5
- Simplified vision model initialization
"""

import os
import re
import io
import fitz
import pytesseract
from PIL import Image
import torch
from transformers import AutoImageProcessor, TableTransformerForObjectDetection
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
from groq import Groq
import warnings
from docx import Document
from pptx import Presentation
import subprocess

# Import our new vision model factory
from vision_models import VisionModelFactory

warnings.filterwarnings("ignore")
os.environ["TOKENIZERS_PARALLELISM"] = "false"


# ... [Keep all the ChartDetector classes exactly as they were] ...
class ChartDetector:
    """Base class for chart detection models."""

    def __init__(self, *args, **kwargs):
        self.device = (
            "mps"
            if torch.backends.mps.is_available()
            else "cuda" if torch.cuda.is_available() else "cpu"
        )
        print(f"ChartDetector using device: {self.device}")
        self.model = self.load_model(*args, **kwargs)

    def load_model(self, *args, **kwargs):
        raise NotImplementedError("Subclasses must implement this method.")

    def detect(self, page_image: Image.Image) -> List[Tuple[int, int, int, int]]:
        raise NotImplementedError("Subclasses must implement this method.")


class HeuristicDetector(ChartDetector):
    """Heuristic-based chart detection logic"""

    def load_model(self, *args, **kwargs):
        print("Using HeuristicDetector.")
        return None

    def detect(self, page_image: Image.Image) -> List[Tuple[int, int, int, int]]:
        boxes = []
        page_width, page_height = page_image.size
        grid_divisions = 3
        section_width = page_width // grid_divisions
        section_height = page_height // grid_divisions

        for row in range(grid_divisions):
            for col in range(grid_divisions):
                x1 = col * section_width
                y1 = row * section_height
                x2 = (col + 1) * section_width
                y2 = (row + 1) * section_height
                section_image = page_image.crop((x1, y1, x2, y2))

                if self._looks_like_chart(section_image):
                    boxes.append((x1, y1, x2, y2))
                    print(
                        f"    Heuristic found potential chart in section ({row},{col})"
                    )
        return boxes

    def _looks_like_chart(self, image: Image.Image) -> bool:
        try:
            width, height = image.size
            if width < 50 or height < 50:
                return False
            ocr_text = pytesseract.image_to_string(image).strip()
            text_density = (len(ocr_text) / (width * height)) * 1000
            if text_density > 0.4:
                return False
            img_array = np.array(image.convert("L"))
            variance = np.var(img_array)
            return variance > 500
        except:
            return False


class TATRDetector(ChartDetector):
    """Table Transformer (TATR) detection"""

    def load_model(self, model_path: str = "microsoft/table-transformer-detection"):
        try:
            print(f"Loading Table Transformer model: {model_path}")
            self.image_processor = AutoImageProcessor.from_pretrained(
                model_path, use_fast=True
            )
            model = TableTransformerForObjectDetection.from_pretrained(model_path)
            model.to(self.device)
            return model
        except Exception as e:
            print(f"Error loading TATR model: {e}")
            return None

    def detect(self, page_image: Image.Image) -> List[Tuple[int, int, int, int]]:
        if not self.model:
            return []
        inputs = self.image_processor(images=page_image, return_tensors="pt").to(
            self.device
        )
        outputs = self.model(**inputs)
        target_sizes = torch.tensor([page_image.size[::-1]]).to(self.device)
        results = self.image_processor.post_process_object_detection(
            outputs, threshold=0.8, target_sizes=target_sizes
        )[0]
        boxes = []
        for score, label, box in zip(
            results["scores"], results["labels"], results["boxes"]
        ):
            box = [round(i, 2) for i in box.tolist()]
            boxes.append(tuple(map(int, box)))
        print(f"    Table Transformer detected {len(boxes)} potential charts/tables.")
        return boxes


@dataclass
class Chunk:
    """Represents a text chunk with metadata"""

    text: str
    source: str
    page: int
    chunk_id: int


class SmartRAG:
    def __init__(
        self,
        model_name: str = "sentence-transformers/all-MiniLM-L6-v2",
        output_dir: str = "potential_charts",
        vision_model_name: str = "Moondream2",  # NEW PARAMETER
    ):
        """
        Initialize the RAG system.

        Args:
            model_name: Sentence transformer model for embeddings
            output_dir: Directory to save extracted charts
            vision_model_name: Vision model to use (Moondream2, Qwen3-VL-2B, InternVL3.5-1B)
        """
        print("Initializing Smart RAG system...")
        print(f"Selected vision model: {vision_model_name}")

        self.groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY", ""))
        self.embedding_model = SentenceTransformer(model_name)
        self.embedding_dim = self.embedding_model.get_sentence_embedding_dimension()

        # Create output directory
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        print(f"Chart output directory: {self.output_dir}")

        # Initialize chart detectors
        self.chart_detector_tatr = TATRDetector()
        self.chart_detector_heuristic = HeuristicDetector()

        # NEW: Load vision model using factory
        self.vision_model_name = vision_model_name
        self.vision_model = None
        self.use_vision = False

        try:
            print(f"\nLoading vision model: {vision_model_name}...")
            self.vision_model = VisionModelFactory.create_model(vision_model_name)

            if self.vision_model:
                self.use_vision = True
                print(f"✓ Vision model '{vision_model_name}' loaded successfully")
            else:
                print(f"✗ Failed to load vision model '{vision_model_name}'")
                print("Continuing without vision capabilities")

        except Exception as e:
            print(f"Error loading vision model: {e}")
            print("Continuing without vision capabilities")
            self.use_vision = False

        # Vector store
        self.index = None
        self.chunks: List[Chunk] = []
        self.chart_descriptions = {}

    # ... [Keep all the utility methods exactly as they were] ...
    def _expand_box(
        self,
        box: Tuple[int, int, int, int],
        page_width: int,
        page_height: int,
        padding_percent: float = 0.1,
    ) -> Tuple[int, int, int, int]:
        x1, y1, x2, y2 = box
        width = x2 - x1
        height = y2 - y1
        pad_x = int(width * padding_percent * 5)
        pad_y = int(height * padding_percent)
        x1 = max(0, x1 - pad_x)
        y1 = max(0, y1 - pad_y)
        x2 = min(page_width, x2 + pad_x)
        y2 = min(page_height, y2 + pad_y)
        return (x1, y1, x2, y2)

    def _calculate_iou(
        self, box1: Tuple[int, int, int, int], box2: Tuple[int, int, int, int]
    ) -> float:
        x1_inter = max(box1[0], box2[0])
        y1_inter = max(box1[1], box2[1])
        x2_inter = min(box1[2], box2[2])
        y2_inter = min(box1[3], box2[3])
        if x2_inter < x1_inter or y2_inter < y1_inter:
            return 0.0
        intersection_area = (x2_inter - x1_inter) * (y2_inter - y1_inter)
        box1_area = (box1[2] - box1[0]) * (box1[3] - box1[1])
        box2_area = (box2[2] - box2[0]) * (box2[3] - box2[1])
        union_area = box1_area + box2_area - intersection_area
        if union_area == 0:
            return 0.0
        return intersection_area / union_area

    def _combine_and_deduplicate_boxes(
        self,
        boxes1: List[Tuple[int, int, int, int]],
        boxes2: List[Tuple[int, int, int, int]],
        iou_threshold: float = 0.5,
    ) -> List[Tuple[int, int, int, int]]:
        if not boxes1 and not boxes2:
            return []
        all_boxes = list(boxes1) + list(boxes2)
        if len(all_boxes) == 0:
            return []
        unique_boxes = []
        for box in all_boxes:
            is_duplicate = False
            for unique_box in unique_boxes:
                iou = self._calculate_iou(box, unique_box)
                if iou > iou_threshold:
                    is_duplicate = True
                    box_area = (box[2] - box[0]) * (box[3] - box[1])
                    unique_box_area = (unique_box[2] - unique_box[0]) * (
                        unique_box[3] - unique_box[1]
                    )
                    if box_area > unique_box_area:
                        idx = unique_boxes.index(unique_box)
                        unique_boxes[idx] = box
                    break
            if not is_duplicate:
                unique_boxes.append(box)
        return unique_boxes

    # UPDATED: Simplified image description methods
    def _describe_image(self, image: Image.Image) -> str:
        """Use selected vision model to describe charts and diagrams"""
        if not self.use_vision or not self.vision_model:
            return "[Vision model not available]"

        prompt = """
Describe this chart or diagram in complete detail. Include ALL data points, trends, labels, axis information, and key insights. List specific numbers, percentages, years, and values shown. Describe the complete trend from beginning to end.
"""
        return self.vision_model.describe_image(image, prompt)

    def _describe_slide(self, slide_image: Image.Image) -> str:
        """Use vision model to provide comprehensive slide description"""
        if not self.use_vision or not self.vision_model:
            return "[Vision model not available]"

        prompt = """
Describe the PowerPoint slide clearly and precisely. Include:

Title/Heading: State what the slide is about.

Text Content: Summarize all visible text (titles, labels, bullets, captions).

Visual Elements: Describe all charts, graphs, diagrams, images, or icons. For charts, include:
- Chart type
- What each axis or category represents
- Key trends, patterns, and directional changes (e.g., increasing, decreasing, flat, spikes, clusters)
- Notable data points or ranges if visible

Layout & Design: Briefly note how the slide is structured and how elements are arranged.

Key Insight: State the main takeaway in one clear sentence.

Be specific with numbers, dates, and labels when they appear. Keep the description factual and concise.
"""
        return self.vision_model.describe_image(slide_image, prompt)

    # ... [Keep ALL remaining methods exactly as they were in the original code] ...
    # This includes: extract_text_from_pptx, _convert_pptx_to_images, extract_text_from_docx,
    # diagnose_pptx, _render_slide_to_image, extract_charts_as_images, _looks_like_chart,
    # _has_chart_aspect_ratio, extract_text_from_pdf, _is_likely_chart, _convert_to_markdown,
    # smart_chunk, _get_page_from_context, index_document, search, query

    def extract_text_from_pptx(
        self, pptx_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        print(f"Processing PPTX: {pptx_path}")
        prs = Presentation(pptx_path)
        full_text = []
        total_slides = len(prs.slides)
        slide_images = self._convert_pptx_to_images(pptx_path, file_output_dir)

        for slide_idx, slide in enumerate(prs.slides):
            print(f"\nProcessing slide {slide_idx + 1}/{total_slides}...")
            if progress_callback:
                progress_callback(slide_idx + 1, total_slides)

            slide_text = [f"## Slide {slide_idx + 1}\n"]
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())

            if text_content:
                slide_text.append("Text content:\n" + "\n".join(text_content) + "\n")

            slide_image = None
            if slide_idx < len(slide_images):
                slide_image = slide_images[slide_idx]

            if slide_image and self.use_vision:
                print(f"  Analyzing full slide with vision model...")
                slide_filename = f"slide{slide_idx + 1}_full.png"
                slide_path = os.path.join(file_output_dir, slide_filename)
                slide_image.save(slide_path)
                description = self._describe_slide(slide_image)
                slide_text.append(f"\n[SLIDE VISUAL DESCRIPTION: {description}]\n")
                self.chart_descriptions[slide_filename] = description

            full_text.append("\n".join(slide_text))

        return "\n\n".join(full_text)

    def _convert_pptx_to_images(
        self, pptx_path: str, output_dir: str
    ) -> List[Image.Image]:
        images = []
        try:
            import tempfile

            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = os.path.join(tmpdir, "presentation.pdf")
                cmd = [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmpdir,
                    pptx_path,
                ]
                result = subprocess.run(
                    cmd, capture_output=True, timeout=120, text=True
                )
                if result.returncode == 0:
                    pdf_files = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]
                    if pdf_files:
                        pdf_path = os.path.join(tmpdir, pdf_files[0])
                        doc = fitz.open(pdf_path)
                        for page_num in range(len(doc)):
                            page = doc[page_num]
                            mat = fitz.Matrix(2.0, 2.0)
                            pix = page.get_pixmap(matrix=mat)
                            img_bytes = pix.tobytes("png")
                            img = Image.open(io.BytesIO(img_bytes))
                            images.append(img)
                        doc.close()
        except Exception as e:
            print(f"  Error converting PPTX to images: {e}")
        return images

    def extract_text_from_docx(
        self, docx_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        print(f"Processing DOCX: {docx_path}")
        doc = Document(docx_path)
        full_text = []

        for para_idx, para in enumerate(doc.paragraphs):
            if para.text.strip():
                full_text.append(para.text)
            if progress_callback and para_idx % 10 == 0:
                progress_callback(para_idx, len(doc.paragraphs))

        for rel_idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image = Image.open(io.BytesIO(image_data))
                    if self._is_likely_chart(image):
                        if self.use_vision:
                            description = self._describe_image(image)
                            full_text.append(f"\n[CHART DESCRIPTION: {description}]\n")
                            chart_filename = f"docx_chart{rel_idx + 1}.png"
                            chart_path = os.path.join(file_output_dir, chart_filename)
                            image.save(chart_path)
                            self.chart_descriptions[chart_filename] = description
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        full_text.append(f"\n[OCR from image: {ocr_text}]\n")
                except Exception as e:
                    print(f"  Error processing image {rel_idx}: {e}")

        return "\n\n".join(full_text)

    def extract_charts_as_images(
        self, page, file_output_dir: str, page_num: int
    ) -> List[Image.Image]:
        charts = []
        try:
            mat = fitz.Matrix(4.0, 4.0)
            pix = page.get_pixmap(matrix=mat)
            page_image = Image.open(io.BytesIO(pix.tobytes("png")))

            chart_boxes_tatr = self.chart_detector_tatr.detect(page_image)
            chart_boxes_heuristic = self.chart_detector_heuristic.detect(page_image)
            chart_boxes_combined = self._combine_and_deduplicate_boxes(
                chart_boxes_tatr, chart_boxes_heuristic, iou_threshold=0.5
            )

            for i, box in enumerate(chart_boxes_combined):
                box = self._expand_box(box, page_image.width, page_image.height)
                chart_image = page_image.crop(box)
                if chart_image.width < 150 or chart_image.height < 150:
                    continue
                charts.append(chart_image)
                chart_filename = f"page{page_num}_chart{i + 1}.png"
                chart_path = os.path.join(file_output_dir, chart_filename)
                chart_image.save(chart_path)

        except Exception as e:
            print(f"    Error during chart detection on page {page_num}: {e}")
        return charts

    def _looks_like_chart(self, image: Image.Image) -> bool:
        try:
            width, height = image.size
            ocr_text = pytesseract.image_to_string(image).strip()
            text_density = (len(ocr_text) / (width * height)) * 1000
            if text_density > 0.4:
                return False
            img_array = np.array(image.convert("L"))
            variance = np.var(img_array)
            return variance > 500
        except:
            img_array = np.array(image.convert("L"))
            variance = np.var(img_array)
            return variance > 500

    def extract_text_from_pdf(
        self, pdf_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        print(f"Processing PDF: {pdf_path}")
        doc = fitz.open(pdf_path)
        full_text = []
        total_pages = len(doc)

        for page_num, page in enumerate(doc):
            print(f"\nProcessing page {page_num + 1}/{len(doc)}...")
            if progress_callback:
                progress_callback(page_num + 1, total_pages)

            page_text = []
            text = page.get_text()
            if text.strip():
                page_text.append(text)

            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))

                    if self._is_likely_chart(image):
                        if self.use_vision:
                            description = self._describe_image(image)
                            page_text.append(f"\n[CHART DESCRIPTION: {description}]\n")
                            chart_filename = (
                                f"page{page_num + 1}_embedded{img_idx + 1}.png"
                            )
                            chart_path = os.path.join(file_output_dir, chart_filename)
                            image.save(chart_path)
                            self.chart_descriptions[chart_filename] = description

                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        page_text.append(f"\n[OCR from image: {ocr_text}]\n")

                except Exception as e:
                    print(f"    Error processing embedded image {img_idx}: {str(e)}")

            chart_images = self.extract_charts_as_images(
                page=page, page_num=page_num + 1, file_output_dir=file_output_dir
            )

            for chart_idx, chart_image in enumerate(chart_images):
                try:
                    if self.use_vision:
                        description = self._describe_image(chart_image)
                        page_text.append(f"\n[CHART DESCRIPTION: {description}]\n")
                        chart_filename = f"page{page_num + 1}_chart{chart_idx + 1}.png"
                        self.chart_descriptions[chart_filename] = description

                    ocr_text = pytesseract.image_to_string(chart_image)
                    if ocr_text.strip():
                        page_text.append(f"\n[OCR from chart region: {ocr_text}]\n")

                except Exception as e:
                    print(f"    Error analyzing chart {chart_idx + 1}: {str(e)}")

            page_markdown = self._convert_to_markdown(
                "\n".join(page_text), page_num + 1
            )
            full_text.append(page_markdown)

        doc.close()
        return "\n\n".join(full_text)

    def _is_likely_chart(self, image: Image.Image) -> bool:
        width, height = image.size
        if width < 200 or height < 200:
            return False
        try:
            ocr_text = pytesseract.image_to_string(image).strip()
            text_density = (
                (len(ocr_text) / (width * height)) * 1000 if (width * height) > 0 else 0
            )
            if text_density > 0.7:
                return False
            return True
        except:
            return True

    def _convert_to_markdown(self, text: str, page_num: int) -> str:
        markdown = f"## Page {page_num}\n\n"
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" +", " ", text)
        lines = text.split("\n")
        formatted_lines = []
        for line in lines:
            line = line.strip()
            if not line:
                formatted_lines.append("")
                continue
            if len(line) < 100 and line.isupper() and len(line) > 3:
                formatted_lines.append(f"### {line.title()}\n")
            else:
                formatted_lines.append(line)
        markdown += "\n".join(formatted_lines)
        return markdown

    def smart_chunk(
        self, text: str, source: str, chunk_size: int = 500, overlap: int = 100
    ) -> List[Chunk]:
        chunks = []
        if "## Slide" in text and source.lower().endswith(".pptx"):
            slides = re.split(r"(## Slide \d+)", text)
            current_slide_num = 0
            for i in range(1, len(slides), 2):
                if i + 1 < len(slides):
                    slide_header = slides[i]
                    slide_content = slides[i + 1].strip()
                    slide_match = re.search(r"## Slide (\d+)", slide_header)
                    if slide_match:
                        current_slide_num = int(slide_match.group(1))
                    if slide_content:
                        full_slide = f"{slide_header}\n{slide_content}"
                        chunks.append(
                            Chunk(
                                text=full_slide,
                                source=source,
                                page=current_slide_num,
                                chunk_id=len(chunks),
                            )
                        )
            return chunks

        text_and_charts = re.split(
            r"(\[CHART DESCRIPTION:.*?\]|\[SLIDE VISUAL DESCRIPTION:.*?\])",
            text,
            flags=re.DOTALL,
        )
        for piece in text_and_charts:
            if piece.startswith("[CHART DESCRIPTION:") or piece.startswith(
                "[SLIDE VISUAL DESCRIPTION:"
            ):
                chunks.append(
                    Chunk(
                        text=piece,
                        source=source,
                        page=self._get_page_from_context(text, piece),
                        chunk_id=len(chunks),
                    )
                )
            else:
                pages = re.split(r"## (?:Page|Slide) (\d+)", piece)
                current_page = 1
                for i in range(1, len(pages), 2):
                    if i < len(pages):
                        current_page = int(pages[i])
                        page_text = pages[i + 1] if i + 1 < len(pages) else ""
                    else:
                        page_text = pages[i]
                    if not page_text.strip():
                        continue
                    sentences = re.split(r"(?<=[.!?])\s+", page_text)
                    current_chunk = []
                    current_length = 0
                    for sentence in sentences:
                        sentence = sentence.strip()
                        if not sentence:
                            continue
                        sentence_length = len(sentence)
                        if (
                            current_length + sentence_length > chunk_size
                            and current_chunk
                        ):
                            chunk_text = " ".join(current_chunk)
                            chunks.append(
                                Chunk(
                                    text=chunk_text,
                                    source=source,
                                    page=current_page,
                                    chunk_id=len(chunks),
                                )
                            )
                            overlap_sentences = []
                            overlap_length = 0
                            for s in reversed(current_chunk):
                                if overlap_length + len(s) <= overlap:
                                    overlap_sentences.insert(0, s)
                                    overlap_length += len(s)
                                else:
                                    break
                            current_chunk = overlap_sentences
                            current_length = overlap_length
                        current_chunk.append(sentence)
                        current_length += sentence_length
                    if current_chunk:
                        chunk_text = " ".join(current_chunk)
                        chunks.append(
                            Chunk(
                                text=chunk_text,
                                source=source,
                                page=current_page,
                                chunk_id=len(chunks),
                            )
                        )
        return chunks

    def _get_page_from_context(self, full_text, chart_description):
        try:
            preceding_text = full_text.split(chart_description)[0]
            page_matches = re.findall(r"## Page (\d+)", preceding_text)
            if page_matches:
                return int(page_matches[-1])
        except:
            pass
        return 0

    def index_document(
        self,
        file_path: str,
        chunk_size: int = 500,
        overlap: int = 100,
        progress_callback=None,
    ):
        file_output_dir = self.output_dir
        file_ext = os.path.splitext(file_path)[1].lower()

        if file_ext == ".pdf":
            markdown_text = self.extract_text_from_pdf(
                file_path, file_output_dir, progress_callback
            )
        elif file_ext == ".docx":
            markdown_text = self.extract_text_from_docx(
                file_path, file_output_dir, progress_callback
            )
        elif file_ext == ".pptx":
            markdown_text = self.extract_text_from_pptx(
                file_path, file_output_dir, progress_callback
            )
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

        chunks = self.smart_chunk(markdown_text, file_path, chunk_size, overlap)
        self.chunks.extend(chunks)
        if not chunks:
            return

        chunk_texts = [chunk.text for chunk in chunks]
        embeddings = self.embedding_model.encode(chunk_texts, show_progress_bar=True)
        if len(embeddings.shape) == 1:
            embeddings = embeddings.reshape(1, -1)
        if self.index is None:
            self.index = faiss.IndexFlatL2(self.embedding_dim)
        self.index.add(np.array(embeddings).astype("float32"))

    def search(self, query: str, top_k: int = 5) -> List[Tuple[Chunk, float]]:
        if self.index is None or len(self.chunks) == 0:
            return []
        query_embedding = self.embedding_model.encode([query])
        distances, indices = self.index.search(
            np.array(query_embedding).astype("float32"), top_k
        )
        results = []
        for dist, idx in zip(distances[0], indices[0]):
            if idx < len(self.chunks):
                results.append((self.chunks[idx], float(dist)))
        return results

    def query(self, question: str, top_k: int = 5) -> Dict:
        results = self.search(question, top_k)
        context = "\n\n---\n\n".join(
            [f"[Page: {chunk.page}]\n{chunk.text}" for chunk, score in results]
        )
        prompt = f"""Based on the following context from a document, answer the question accurately and concisely.
        Context:
        {context}
        Question: {question}
        Instructions:
        - Answer directly based on the context provided
        - If chart data is mentioned, cite specific values and trends
        - If the answer isn't in the context, say so
        - Keep your answer focused and under 150 words
        Answer:"""
        try:
            response = self.groq_client.chat.completions.create(
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                messages=[
                    {
                        "role": "system",
                        "content": "You are a precise assistant that answers questions based strictly on the provided context. When chart data is available, cite specific numbers and trends. Be concise and accurate. If page numbers are in the metadata, mention them.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.1,
                max_tokens=500,
            )
            return {
                "question": question,
                "context": context,
                "results": [
                    {
                        "text": chunk.text,
                        "source": chunk.source,
                        "page": chunk.page,
                        "score": score,
                    }
                    for chunk, score in results
                ],
                "response": response.choices[0].message.content,
            }
        except Exception as e:
            return {"error": str(e)}
