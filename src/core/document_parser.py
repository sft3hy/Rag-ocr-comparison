import os
import re
import io
import fitz  # PyMuPDF
import pytesseract
import subprocess
from PIL import Image
from docx import Document
from pptx import Presentation
from typing import List, Tuple, Dict, Optional
import numpy as np

# Imports from your architecture
from src.vision.vision_models import VisionModel
from src.utils.chart_detection import PubLayNetDetector


class DocumentParser:
    """
    Handles the parsing of different document formats (PDF, DOCX, PPTX).

    Strategy:
    1. Convert pages to images.
    2. Use PubLayNet (Detectron2) to Detect & Crop charts/figures.
    3. Use Vision Model to Describe the cropped visuals.
    """

    def __init__(self, vision_model: Optional[VisionModel], output_dir: str):
        """
        Initializes the parser.

        Args:
            vision_model (VisionModel, optional): An instance of a loaded vision model.
            output_dir (str): Directory to save extracted images.
        """
        self.vision_model = vision_model
        self.use_vision = vision_model is not None
        self.output_dir = output_dir
        self.chart_descriptions: Dict[str, str] = {}

        # Initialize the Layout Detector
        self.layout_detector = PubLayNetDetector(confidence_threshold=0.5, padding=60)

    def parse(self, file_path: str, progress_callback=None) -> str:
        """
        Public method to parse a document based on its file extension.
        """
        file_ext = os.path.splitext(file_path)[1].lower()

        # Create file-specific output directory
        file_output_dir = os.path.join(
            self.output_dir, os.path.splitext(os.path.basename(file_path))[0]
        )
        os.makedirs(file_output_dir, exist_ok=True)

        try:
            if file_ext == ".pdf":
                result = self._extract_text_from_pdf(
                    file_path, file_output_dir, progress_callback
                )
            elif file_ext == ".docx":
                result = self._extract_text_from_docx(
                    file_path, file_output_dir, progress_callback
                )
            elif file_ext == ".pptx":
                result = self._extract_text_from_pptx(
                    file_path, file_output_dir, progress_callback
                )
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")

            return result

        finally:
            # Optional: Offload detector to free GPU memory after parsing is done
            # self.layout_detector.offload_model()
            pass

    def _process_page_visuals(
        self, page_image: Image.Image, page_id: str, output_dir: str
    ) -> str:
        """
        1. Detects charts/figures using PubLayNet.
        2. Crops them.
        3. Sends crops to Vision Model for description.
        """
        if not self.use_vision:
            return ""

        # 1. Detect
        bboxes = self.layout_detector.detect(page_image)

        if not bboxes:
            return ""

        descriptions = []
        print(f"  > Found {len(bboxes)} visual elements on {page_id}")

        # 2. Iterate, Crop, Describe
        for i, (x1, y1, x2, y2) in enumerate(bboxes):
            # Crop
            crop = page_image.crop((x1, y1, x2, y2))

            # Save Crop
            chart_filename = f"{page_id}_visual_{i+1}.png"
            chart_path = os.path.join(output_dir, chart_filename)
            crop.save(chart_path)

            # Vision Prompt
            prompt = (
                "Analyze this specific image crop from a document. "
                "Identify if it is a Chart, Table, or Diagram. "
                "Provide a detailed description of the data, including title, axis labels, "
                "key trends, and specific data points shown."
            )

            # Describe
            try:
                desc = self.vision_model.describe_image(crop, prompt)

                # Filter out empty/failure responses
                if desc and "error" not in desc.lower():
                    formatted_desc = (
                        f"\n> **[Visual Element {i+1}]:**\n"
                        f"> *Image saved to: {chart_filename}*\n"
                        f"> {desc}\n"
                    )
                    descriptions.append(formatted_desc)
                    self.chart_descriptions[chart_filename] = desc
            except Exception as e:
                print(f"    Error describing visual {i+1}: {e}")

        return "\n".join(descriptions)

    def _extract_text_from_pdf(
        self, pdf_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        print(f"Processing PDF: {pdf_path}")
        doc = fitz.open(pdf_path)
        full_text = []
        total_pages = len(doc)

        for page_num, page in enumerate(doc):
            print(f"\nProcessing page {page_num + 1}/{total_pages}...")
            if progress_callback:
                progress_callback(page_num + 1, total_pages)

            # 1. Extract raw text
            page_text = []
            text = page.get_text()
            if text.strip():
                page_text.append(text)

            # 2. Render Page for Visual Analysis
            # Matrix=2.0 for higher resolution (better detection)
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            page_image = Image.open(io.BytesIO(pix.tobytes("png")))

            # 3. Detect & Describe Charts (PubLayNet + Vision)
            if self.use_vision:
                visual_context = self._process_page_visuals(
                    page_image, f"page{page_num + 1}", file_output_dir
                )

                if visual_context:
                    page_text.append(
                        f"\n--- Visual Analysis ---\n{visual_context}\n-----------------------\n"
                    )

            # 4. Format
            page_markdown = self._convert_to_markdown(
                "\n".join(page_text), page_num + 1
            )
            full_text.append(page_markdown)

        doc.close()
        return "\n\n".join(full_text)

    def _extract_text_from_pptx(
        self, pptx_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        print(f"Processing PPTX: {pptx_path}")
        prs = Presentation(pptx_path)
        full_text = []
        total_slides = len(prs.slides)

        # Convert all slides to images first (needed for layout detection)
        slide_images = self._convert_pptx_to_images(pptx_path, file_output_dir)

        for slide_idx, slide in enumerate(prs.slides):
            print(f"\nProcessing slide {slide_idx + 1}/{total_slides}...")
            if progress_callback:
                progress_callback(slide_idx + 1, total_slides)

            slide_text = [f"## Slide {slide_idx + 1}\n"]

            # Extract text from shapes
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())

            if text_content:
                slide_text.append("Text content:\n" + "\n".join(text_content) + "\n")

            # Visual analysis
            if slide_idx < len(slide_images) and self.use_vision:
                slide_image = slide_images[slide_idx]

                visual_context = self._process_page_visuals(
                    slide_image, f"slide{slide_idx + 1}", file_output_dir
                )

                if visual_context:
                    slide_text.append(
                        f"\n--- Visual Analysis ---\n{visual_context}\n-----------------------\n"
                    )

            full_text.append("\n".join(slide_text))

        return "\n\n".join(full_text)

    def _extract_text_from_docx(
        self, docx_path: str, file_output_dir: str, progress_callback=None
    ) -> str:
        """
        DOCX Strategy:
        Layout detection is hard on DOCX without converting to PDF first.
        We stick to extracting embedded images (OLE/Shapes) and running VLM on those.
        """
        print(f"Processing DOCX: {docx_path}")
        doc = Document(docx_path)
        full_text = []

        for para_idx, para in enumerate(doc.paragraphs):
            if para.text.strip():
                full_text.append(para.text)
            if progress_callback and para_idx % 10 == 0:
                progress_callback(para_idx, len(doc.paragraphs))

        # Iterate over relationships to find images
        for rel_idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image = Image.open(io.BytesIO(image_data))

                    # Basic size check to skip icons/lines
                    if image.width > 150 and image.height > 150:
                        if self.use_vision:
                            print(f"  Analyzing embedded image {rel_idx}...")

                            # Save
                            chart_filename = f"docx_image{rel_idx + 1}.png"
                            chart_path = os.path.join(file_output_dir, chart_filename)
                            image.save(chart_path)

                            # Describe directly (no need for PubLayNet here, we already have the crop)
                            prompt = "Analyze this image. Identify if it is a Chart or Table and describe the data trends."
                            desc = self.vision_model.describe_image(image, prompt)

                            if desc:
                                full_text.append(
                                    f"\n> **[Embedded Image {rel_idx + 1}]:**\n> {desc}\n"
                                )
                                self.chart_descriptions[chart_filename] = desc

                except Exception as e:
                    print(f"  Error processing image {rel_idx}: {e}")

        return "\n\n".join(full_text)

    def _convert_pptx_to_images(
        self, pptx_path: str, output_dir: str
    ) -> List[Image.Image]:
        """Converts PPTX slides to images via PDF intermediate."""
        images = []
        try:
            import tempfile

            # Create a dedicated temp dir for conversion
            with tempfile.TemporaryDirectory() as tmpdir:
                # Requires LibreOffice (soffice)
                cmd = [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmpdir,
                    pptx_path,
                ]

                # Run with timeout to prevent hanging
                result = subprocess.run(
                    cmd, capture_output=True, timeout=120, text=True
                )

                # Find the generated PDF
                pdf_files = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]

                if result.returncode == 0 and pdf_files:
                    pdf_path = os.path.join(tmpdir, pdf_files[0])
                    doc = fitz.open(pdf_path)

                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        mat = fitz.Matrix(2.0, 2.0)
                        pix = page.get_pixmap(matrix=mat)
                        img_bytes = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_bytes))
                        images.append(img)
                    doc.close()
                else:
                    print(
                        "  Warning: soffice conversion failed. Is LibreOffice installed?"
                    )
                    print(f"  Stderr: {result.stderr}")
        except Exception as e:
            print(f"  Error converting PPTX to images: {e}")
        return images

    def _convert_to_markdown(self, text: str, page_num: int) -> str:
        markdown = f"## Page {page_num}\n\n"
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" +", " ", text)
        lines = text.split("\n")
        formatted_lines = []
        for line in lines:
            line = line.strip()
            if not line:
                formatted_lines.append("")
                continue
            if len(line) < 100 and line.isupper() and len(line) > 3:
                formatted_lines.append(f"### {line.title()}\n")
            else:
                formatted_lines.append(line)
        markdown += "\n".join(formatted_lines)
        return markdown
